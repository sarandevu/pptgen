import streamlit as st
import json
import io
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfReader

# ---------------------------
# 1. Gemini Client Setup
# ---------------------------
client = genai.Client(api_key=st.secrets["GEMINI_API_KEY"])

# ---------------------------
# Helper: Safe JSON parsing
# ---------------------------
def safe_json(text):
    try:
        return json.loads(text)
    except:
        start = text.find("{")
        end = text.rfind("}") + 1
        return json.loads(text[start:end])

# ---------------------------
# 2. Knowledge Extraction Agent
# ---------------------------
def extract_knowledge(text):

    prompt = f"""
Extract factual knowledge from this content.

Return JSON:

{{
 "topic":"",
 "key_concepts":[],
 "definitions":[],
 "important_facts":[],
 "formulas":[],
 "examples":[]
}}

Text:
{text[:8000]}
"""

    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt,
        config=types.GenerateContentConfig(response_mime_type="application/json")
    )

    return safe_json(response.text)


# ---------------------------
# 3. Presentation Architect Agent
# ---------------------------
def generate_presentation(notes, instructions, image_file=None):

    system = """
You are a Professional Technical Presentation Architect.

Convert notes into structured slides.

Structure:
Title
Problem
Existing Solutions
Proposed System
Architecture
Key Components
Advantages
Conclusion

Return JSON:
{
 "title":"",
 "image_role":"background | logo | none",
 "slides":[
  {
   "title":"",
   "bullets":[],
   "speaker_notes":""
  }
 ]
}
"""

    contents = [
        f"Notes: {notes}",
        f"Design: {instructions}"
    ]

    if image_file:
        image_file.seek(0)
        contents.append(
            types.Part.from_bytes(
                data=image_file.read(),
                mime_type=image_file.type
            )
        )

    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=contents,
        config=types.GenerateContentConfig(
            system_instruction=system,
            response_mime_type="application/json"
        )
    )

    return safe_json(response.text)


# ---------------------------
# 4. Slide Improver Agent
# ---------------------------
def improve_slide(text):

    prompt = f"""
Improve this PowerPoint slide.

Make it:
• concise
• technical
• professional

Slide:
{text}
"""

    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt
    )

    return response.text


# ---------------------------
# 5. Design Advisor Agent
# ---------------------------
def get_design_suggestions(topic):

    prompt = f"""
Suggest professional PowerPoint design ideas.

Topic: {topic}

Return JSON:

{{
 "theme_colors":"",
 "background_style":"",
 "font_style":"",
 "visual_elements":[]
}}
"""

    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt,
        config=types.GenerateContentConfig(response_mime_type="application/json")
    )

    return safe_json(response.text)


# ---------------------------
# 6. Research Mode
# ---------------------------
def extract_pdf_text(pdf):

    reader = PdfReader(pdf)

    text = ""

    for page in reader.pages:
        t = page.extract_text()
        if t:
            text += t

    return text


def research_to_slides(text):

    prompt = f"""
Convert research text into slides.

Structure:
Title
Background
Problem
Method
Results
Conclusion

Return JSON presentation.

Text:
{text[:8000]}
"""

    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt,
        config=types.GenerateContentConfig(response_mime_type="application/json")
    )

    return safe_json(response.text)


# ---------------------------
# 7. PPT Builder
# ---------------------------
def build_presentation(data, image_file=None):

    prs = Presentation()

    role = data.get("image_role", "none")

    image_bytes = None

    if image_file:
        image_file.seek(0)
        image_bytes = image_file.read()

    def apply_assets(slide):

        if not image_bytes:
            return

        if role == "background":

            slide.shapes.add_picture(
                io.BytesIO(image_bytes),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height
            )

        elif role == "logo":

            slide.shapes.add_picture(
                io.BytesIO(image_bytes),
                Inches(8.5),
                Inches(0.2),
                width=Inches(1)
            )

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = data.get("title", "Presentation")
    apply_assets(slide)

    # Content slides
    for s in data.get("slides", []):

        slide = prs.slides.add_slide(prs.slide_layouts[1])

        apply_assets(slide)

        slide.shapes.title.text = s.get("title", "")

        tf = slide.placeholders[1].text_frame
        tf.clear()

        for bullet in s.get("bullets", []):
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

        # speaker notes
        notes = s.get("speaker_notes")

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    ppt = io.BytesIO()
    prs.save(ppt)

    return ppt.getvalue()


# ---------------------------
# 8. Streamlit UI
# ---------------------------

st.set_page_config(page_title="AI PPT Architect", layout="wide")

st.title("🏗️ AI Knowledge-Driven PPT Architect")

if "architect_data" not in st.session_state:
    st.session_state.architect_data = None

# Sidebar
with st.sidebar:

    st.header("Create Presentation")

    notes = st.text_area("Project Notes")

    design_prompt = st.text_input("Design instructions")

    image_file = st.file_uploader("Upload optional image", type=["png", "jpg", "jpeg"])

    use_kg = st.checkbox("Use Knowledge Extraction")

    if st.button("Generate PPT Structure"):

        if notes:

            with st.spinner("Processing..."):

                try:

                    processed_notes = notes

                    if use_kg:
                        kg = extract_knowledge(notes)
                        processed_notes = json.dumps(kg, indent=2)
                        st.session_state.extracted_knowledge = kg
                    else:
                        st.session_state.extracted_knowledge = None

                    st.session_state.architect_data = generate_presentation(
                        processed_notes,
                        design_prompt,
                        image_file
                    )

                except Exception as e:

                    st.error(e)

        else:
            st.error("Enter notes first")

    st.divider()

    st.header("Research Mode")

    pdf = st.file_uploader("Upload research PDF", type=["pdf"])

    if pdf and st.button("Generate Slides from PDF"):

        with st.spinner("Analyzing research..."):

            try:

                text = extract_pdf_text(pdf)

                st.session_state.architect_data = research_to_slides(text)

            except Exception as e:

                st.error(e)


# ---------------------------
# 9. Main Workspace
# ---------------------------

if st.session_state.architect_data:

    data = st.session_state.architect_data

    st.header(data.get("title", "Presentation"))

    if st.button("Get Design Suggestions"):

        suggestions = get_design_suggestions(data.get("title"))

        st.write("Theme:", suggestions.get("theme_colors"))
        st.write("Background:", suggestions.get("background_style"))
        st.write("Font:", suggestions.get("font_style"))
        st.write("Visuals:", suggestions.get("visual_elements"))

    st.divider()

    for i, slide in enumerate(data.get("slides", [])):

        with st.expander(f"Slide {i+1}: {slide.get('title')}"):

            text = "\n".join(slide.get("bullets", []))

            edited = st.text_area(
                "Edit slide",
                text,
                key=f"slide_{i}"
            )

            slide["bullets"] = edited.split("\n")

            if st.button(f"Improve Slide {i+1}"):

                improved = improve_slide(edited)

                st.write(improved)

    st.divider()

    if st.checkbox("I verify content is correct"):

        ppt = build_presentation(data, image_file)

        st.download_button(
            "Download PPT",
            ppt,
            file_name="ai_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )